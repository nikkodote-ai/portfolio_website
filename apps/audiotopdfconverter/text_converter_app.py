import os
import re
import sys
import time
import wget
import PyPDF2
from boto3 import Session, resource
from botocore.exceptions import BotoCoreError, ClientError
from pptx import Presentation

# from https://docs.aws.amazon.com/polly/latest/dg/get-started-what-next.html
AWSAccessKeyId = os.getenv('AWSAccessKeyId')
AWSSecretKey = os.getenv('AWSSecretKey')
s3 = resource('s3', region_name="ap-southeast-2", aws_access_key_id=AWSAccessKeyId, aws_secret_access_key=AWSSecretKey)


def convert_to_text(pdf_text):
    # create an object variable in rb mode
    with open(pdf_text, 'rb') as pdf_file:
        # create a pdf reader
        pdf_reader = PyPDF2.PdfFileReader(pdf_file)
        # number of pages in the pdf
        num_pages = pdf_reader.getNumPages()
        # get all the pages from the beginning of the pdf to last
        pages = pdf_reader.getPage(num_pages - 1)
        # store text
        text = pages.extractText()
        text = text.split('\n')
        text = ''.join(text)

        text = re.sub(r"(@\[A-Za-z0-9]+)|([^0-9A-Za-z \t])|(\w+�:\/\/\S+)|^rt|http.+?", "", text)

        with open('converted_to_text.txt', 'w') as output:
            output.writelines(text)
        return text


def convert_ppt_to_text(file_name):
    prs = Presentation(file_name)
    text_runs = {}

    for index, slide in enumerate(prs.slides):
        text_runs[index + 1] = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text_runs[index + 1].append(f"<p>{run.text}</p>")

    final_text = ''
    for slide_index, texts in text_runs.items():
        final_text += f'<break strength ="medium"/> Slide {slide_index}, {" ".join(texts)}'

    final_text = '<speak><prosody rate="85%">' + final_text + '</prosody></speak>'
    # with open("ppt_extracted_text.txt", "w") as output:
    #     output.write(final_text)
    return final_text


def convert_to_audio(text_input, voice_id, engine, file_name):
    # Create a client using the credentials and region defined in the [adminuser]
    # section of the AWS credentials file (~/.aws/credentials).
    session = Session(aws_access_key_id=AWSAccessKeyId, aws_secret_access_key=AWSSecretKey,
                      region_name='ap-southeast-2')
    polly = session.client("polly")

    try:
        # Request speech synthesis
        response = polly.start_speech_synthesis_task(Text=text_input, OutputFormat="mp3",
                                                     VoiceId=voice_id, TextType="ssml", Engine=engine,
                                                     OutputS3BucketName='nikkodoteapps')
        # else:
        #     response = polly.synthesize_speech(Text=text_input, OutputFormat="mp3",
        #                                                  VoiceId=voice_id, TextType = "ssml", Engine=engine)

    except (BotoCoreError, ClientError) as error:
        # The service returned an error, exit gracefully
        print(error)
        sys.exit(-1)

    # Characters less than 3000 can be converted using synthesize_speech
    # and can be saved in the local disk right away
    # if "AudioStream" in response:
    #         with closing(response["AudioStream"]) as stream:
    #            output = f"{file_name}.mp3"
    #            print(output)
    #            try:
    #             # Open a file for writing the output as a binary stream
    #                 with open(output, "wb") as file:
    #                    file.write(stream.read())
    #                 response = wget.download(output)
    #
    #            except IOError as error:
    #               # Could not write to file, exit gracefully
    #               print(error)
    #               sys.exit(-1)

    # Long audio files needs asynchronous synthesis saving the output
    # in the S3 bucket and not locally first then locally
    if 'OutputUri' in response['SynthesisTask']:
        print(response)
        return long_audio_download(response['SynthesisTask']['OutputUri'])
    else:
        # The response didn't contain audio data, exit gracefully
        print("Could not stream audio")
        sys.exit(-1)


def long_audio_download(output_uri):
    output_uri = output_uri.split('/')
    print(output_uri)
    bucket_name = output_uri[-2]
    object_name = output_uri[-1]
    print(f"bucket_name: {bucket_name}")
    print(f"object_name: {object_name}")
    print(bucket_name + "\n" + object_name)
    my_bucket = s3.Bucket(bucket_name)
    files = my_bucket.objects.all()
    bucket_files = [file.key for file in files]
    print(f'bucket files: {bucket_files}')
    print(f'my object name to match : {object_name}')

    while True:
        if object_name not in bucket_files:
            print('File not yet found. Will retry')
            print(f'try going to https://{bucket_name}.s3.amazonaws.com/{object_name} ')
            print(f'>>>bucket files: {bucket_files}')
            print(f'>>>my object name to match : {object_name}')
            try:
                return download_file(bucket_name, object_name)
            except:
                print('Error while trying')
            time.sleep(5)
            continue
        break

    try:
        return download_file(bucket_name, object_name)
    except ClientError as e:
        print(e)
        print("Processing. Download once done")
        time.sleep(25)
        return download_file(bucket_name, object_name)
    else:
        print("Download not done")


def delete_file(bucket_name, file_name):
    my_bucket = s3.Bucket(bucket_name)
    my_bucket.Object(file_name).delete()


def download_file(bucket_name, file_name):
    my_bucket = s3.Bucket(bucket_name)
    file_object = my_bucket.Object(file_name).get()
    return file_object

